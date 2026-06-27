#!/usr/bin/env python3
"""将8张高清截图组装为可编辑的 16:9 PPTX 文件"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# 路径配置
SCREENSHOTS_DIR = '/Users/summer/WorkBuddy/2026-06-27-11-18-55/ppt/screenshots'
OUTPUT_PATH = '/Users/summer/WorkBuddy/2026-06-27-11-18-55/ppt/知几_黑客松路演.pptx'
NUM_SLIDES = 8

# 创建 16:9 演示文稿 (宽13.333英寸 × 高7.5英寸)
prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9 标准宽度
prs.slide_height = Inches(7.5)     # 16:9 标准高度

# 使用空白布局（无默认文本框）
blank_layout = prs.slide_layouts[6]  # Blank layout

for i in range(1, NUM_SLIDES + 1):
    img_path = os.path.join(SCREENSHOTS_DIR, f'slide_{i}.png')
    
    if not os.path.exists(img_path):
        print(f"⚠️  Warning: {img_path} not found, skipping slide {i}")
        continue
    
    print(f"Processing slide {i}...")
    slide = prs.slides.add_slide(blank_layout)
    
    # 将截图作为整页背景图插入
    # 使用 fill 方式：图片铺满整个幻灯片区域
    pic = slide.shapes.add_picture(
        img_path,
        left=Inches(0),
        top=Inches(0),
        width=prs.slide_width,
        height=prs.slide_height
    )

print(f"\n✅ PPTX 已生成: {OUTPUT_PATH}")
print(f"   共 {len(prs.slides)} 页 · 16:9 比例 · 高清背景图")
print(f"   可在 PowerPoint 中添加录屏、文字框、标注等")

prs.save(OUTPUT_PATH)
